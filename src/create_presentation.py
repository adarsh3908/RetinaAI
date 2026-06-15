import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Premium Bright & Modern)
    BG_COLOR = RGBColor(248, 250, 252)       # Light Slate 50 (Very light gray-blue background)
    CARD_BG = RGBColor(255, 255, 255)        # Crisp White
    TEXT_DARK = RGBColor(15, 23, 42)          # Slate 900 (Dark slate for primary text)
    TEXT_MUTED = RGBColor(71, 85, 105)        # Slate 600 (Medium slate for body and notes)
    ACCENT_BLUE = RGBColor(29, 78, 216)       # Blue 700 (Cobalt Blue for structure/primary highlights)
    ACCENT_GREEN = RGBColor(4, 120, 87)       # Emerald 700 (Vibrant green for solutions/positives)
    ACCENT_RED = RGBColor(185, 28, 28)        # Red 700 (Crimson red for problems/warnings)
    BORDER_COLOR = RGBColor(226, 232, 240)    # Slate 200 (Light border gray)
    WHITE = RGBColor(255, 255, 255)

    # Fonts
    FONT_TITLE = "Segoe UI"
    FONT_BODY = "Segoe UI"

    def apply_bright_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_slide_header(slide, title_text, category_text="RETINAAI"):
        # Category/Tracker at top left
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_TITLE
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_COLOR, line_width=1):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(line_width)
        return shape

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide1)

    # Accent decorative block on left
    glow = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = ACCENT_BLUE
    glow.line.fill.background()

    # Dynamic geometric background accent (Circle bottom-right)
    bg_circle = slide1.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(3.5), Inches(4.5), Inches(4.5)
    )
    bg_circle.fill.solid()
    bg_circle.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate 100
    bg_circle.line.fill.background()

    # Title & Subtitle in single text frame
    main_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
    tf = main_box.text_frame
    tf.word_wrap = True

    # Main Title
    p1 = tf.paragraphs[0]
    p1.text = "Student Dropout Risk Prediction"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_DARK
    p1.space_after = Pt(10)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "A Multimodal Machine Learning Approach using Academics, Attendance, and Counselor Insights"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_BLUE
    p2.space_after = Pt(50)

    # Author
    p3 = tf.add_paragraph()
    p3.text = "Presented by: Adarsh Prakash Singh"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_DARK
    p3.space_after = Pt(5)

    # Org
    p4 = tf.add_paragraph()
    p4.text = "RetinaAI System | Kaggle Competition Submission"
    p4.font.name = FONT_BODY
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: Project Overview & Problem Statement
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide2)
    add_slide_header(slide2, "Project Overview & Problem Statement", "Introduction")

    # Left Card: The Problem (Crimson tinted)
    L_RED_BG = RGBColor(254, 242, 242) # Red 50
    add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), bg_color=L_RED_BG, border_color=ACCENT_RED, line_width=2)
    
    # Problem Icon (Circle)
    prob_icon = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), Inches(2.0), Inches(0.4), Inches(0.4))
    prob_icon.fill.solid()
    prob_icon.fill.fore_color.rgb = ACCENT_RED
    prob_icon.line.fill.background()

    prob_box = slide2.shapes.add_textbox(Inches(1.6), Inches(1.95), Inches(4.5), Inches(4.2))
    tf_prob = prob_box.text_frame
    tf_prob.word_wrap = True
    
    p = tf_prob.paragraphs[0]
    p.text = "THE PROBLEM"
    p.font.name = FONT_TITLE
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.space_after = Pt(14)

    bullets_prob = [
        "Educational institutions struggle to identify at-risk students before they drop out.",
        "Key indicators (academic decline, irregular attendance, behavioral changes) are scattered across disconnected systems.",
        "Traditional intervention methods are reactive and occur too late to support the student effectively.",
        "Manual tracking of qualitative insights (e.g., counselor notes) is labor-intensive and rarely integrated with quantitative metrics."
    ]
    for b in bullets_prob:
        p = tf_prob.add_paragraph()
        p.text = "• " + b
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # Right Card: The Solution (Green tinted)
    L_GREEN_BG = RGBColor(240, 253, 250) # Emerald 50
    add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), bg_color=L_GREEN_BG, border_color=ACCENT_GREEN, line_width=2)
    
    # Solution Icon (Circle)
    sol_icon = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.1), Inches(2.0), Inches(0.4), Inches(0.4))
    sol_icon.fill.solid()
    sol_icon.fill.fore_color.rgb = ACCENT_GREEN
    sol_icon.line.fill.background()

    sol_box = slide2.shapes.add_textbox(Inches(7.6), Inches(1.95), Inches(4.5), Inches(4.2))
    tf_sol = sol_box.text_frame
    tf_sol.word_wrap = True

    p = tf_sol.paragraphs[0]
    p.text = "THE RETINAAI SOLUTION"
    p.font.name = FONT_TITLE
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)

    bullets_sol = [
        "A multimodal AI prediction system that fuses academic records, attendance history, and unstructured counselor notes.",
        "Predicts student risk in real-time, categorizing them into: Low Risk (0), Medium Risk (1), or High Risk (2).",
        "NLP extraction translates messy counselor narratives into predictive quantitative dimensions.",
        "Enables targeted, proactive counseling and academic support prior to final withdrawal."
    ]
    for b in bullets_sol:
        p = tf_sol.add_paragraph()
        p.text = "• " + b
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 3: Data Modalities (Multimodal Input)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide3)
    add_slide_header(slide3, "Multimodal Data Modalities", "Data Engineering")

    # 3 Columns for 3 Modalities (White cards with Top Badges)
    modalities = [
        {
            "title": "Tabular Data",
            "desc": "Socio-economic & demographic features.",
            "points": [
                "Demographics: Gender, Branch",
                "Socio-Economic: Family Income, Parent Education, Scholarship Status",
                "Lifestyle Factors: Screen Time, Commute Time, Part-time Job"
            ],
            "accent": ACCENT_BLUE,
            "bg": RGBColor(239, 246, 255) # Blue 50
        },
        {
            "title": "Attendance Time-Series",
            "desc": "Temporal attendance tracking.",
            "points": [
                "Semester-wise attendance rates",
                "Long-term trend indicators",
                "Attendance consistency variance",
                "Linear regression slope metrics"
            ],
            "accent": ACCENT_GREEN,
            "bg": RGBColor(240, 253, 250) # Emerald 50
        },
        {
            "title": "Counsellor Notes (Text)",
            "desc": "Qualitative behavioral insights.",
            "points": [
                "Messy, unstructured counselor logs",
                "Student feedback narratives",
                "Explanatory context for academic decline",
                "Behavioral and emotional observations"
            ],
            "accent": ACCENT_BLUE,
            "bg": RGBColor(239, 246, 255) # Blue 50
        }
    ]

    card_w = Inches(3.64)
    card_h = Inches(4.3)
    gap = Inches(0.4)
    left_margin = Inches(0.8)

    for i, mod in enumerate(modalities):
        x = left_margin + i * (card_w + gap)
        # Background card
        add_card(slide3, x, Inches(2.2), card_w, card_h, bg_color=CARD_BG, border_color=BORDER_COLOR)

        # Circular Badge overlapping at top
        circle = slide3.shapes.add_shape(
            MSO_SHAPE.OVAL, x + card_w/2 - Inches(0.6), Inches(1.5), Inches(1.2), Inches(1.2)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = mod["accent"]
        circle.line.fill.background()
        tf_c = circle.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = str(i+1)
        p_c.alignment = PP_ALIGN.CENTER
        p_c.font.name = FONT_TITLE
        p_c.font.size = Pt(22)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE
        
        box = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.8), card_w - Inches(0.4), card_h - Inches(0.8))
        tf_mod = box.text_frame
        tf_mod.word_wrap = True

        p = tf_mod.paragraphs[0]
        p.text = mod["title"].upper()
        p.font.name = FONT_TITLE
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = mod["accent"]
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

        p = tf_mod.add_paragraph()
        p.text = mod["desc"]
        p.font.name = FONT_BODY
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(14)

        for pt in mod["points"]:
            p = tf_mod.add_paragraph()
            p.text = "• " + pt
            p.font.name = FONT_BODY
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(6)

    # ==========================================
    # SLIDE 4: Feature Engineering Details (Horizontal Process Chevrons)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide4)
    add_slide_header(slide4, "Feature Engineering Strategy", "Data Engineering")

    # 3 Chevrons laying horizontally representing process flow
    feat_groups = [
        {
            "title": "Academic Metrics",
            "accent": ACCENT_BLUE,
            "bg": RGBColor(239, 246, 255), # Blue 50
            "points": [
                "CGPA Mean: Captures overall performance history.",
                "CGPA Standard Deviation: Monitors performance volatility.",
                "CGPA Trend: Captures directional academic trajectory."
            ]
        },
        {
            "title": "Backlog Indicators",
            "accent": TEXT_MUTED,
            "bg": CARD_BG,
            "points": [
                "Total Backlogs: Cumulative failed courses.",
                "Average Backlogs: Course failure rate.",
                "Backlog Trend: Shows whether backlogs are increasing or decreasing over semesters."
            ]
        },
        {
            "title": "Attendance Dynamics",
            "accent": ACCENT_GREEN,
            "bg": RGBColor(240, 253, 250), # Emerald 50
            "points": [
                "Attendance Range & Std Dev: Tracks attendance consistency.",
                "Attendance Slope: Identifies steep drops in class presence.",
                "Semester Metrics: Tracks semester-over-semester differences."
            ]
        }
    ]

    chev_w = Inches(3.7)
    chev_h = Inches(4.3)
    chev_gap = Inches(0.2)
    left_m = Inches(0.8)

    for i, fg in enumerate(feat_groups):
        cx = left_m + i * (chev_w + chev_gap)
        
        # Chevron Shape
        chev = slide4.shapes.add_shape(
            MSO_SHAPE.CHEVRON, cx, Inches(1.8), chev_w, chev_h
        )
        chev.fill.solid()
        chev.fill.fore_color.rgb = fg["bg"]
        chev.line.color.rgb = BORDER_COLOR
        chev.line.width = Pt(1.5)
        
        # Add overlay textbox on top of Chevron for precise alignment
        box = slide4.shapes.add_textbox(cx + Inches(0.5), Inches(2.0), chev_w - Inches(0.8), chev_h - Inches(0.4))
        tf_fg = box.text_frame
        tf_fg.word_wrap = True

        p = tf_fg.paragraphs[0]
        p.text = fg["title"].upper()
        p.font.name = FONT_TITLE
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = fg["accent"]
        p.space_after = Pt(14)

        for pt in fg["points"]:
            p = tf_fg.add_paragraph()
            parts = pt.split(":")
            # Bold the prefix
            run1 = p.add_run()
            run1.text = "• " + parts[0] + ":"
            run1.font.bold = True
            run1.font.name = FONT_BODY
            run1.font.size = Pt(12)
            run1.font.color.rgb = TEXT_DARK
            
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.name = FONT_BODY
            run2.font.size = Pt(12)
            run2.font.color.rgb = TEXT_MUTED
            
            p.space_after = Pt(8)

    # ==========================================
    # SLIDE 5: NLP Processing of Counselor Notes (Flowchart Diagram)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide5)
    add_slide_header(slide5, "NLP Extraction on Counselor Notes", "Text Mining")

    # Draw Flowchart Shapes: Input Document -> Process Hexagon -> Output Circle
    
    # 1. Input Document shape
    add_card(slide5, Inches(0.8), Inches(1.8), Inches(3.2), Inches(3.6), bg_color=CARD_BG, border_color=ACCENT_BLUE)
    doc_title_box = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(2.8), Inches(3.2))
    tf_doc = doc_title_box.text_frame
    tf_doc.word_wrap = True
    p = tf_doc.paragraphs[0]
    p.text = "INPUT: RAW NOTES"
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(14)
    
    p = tf_doc.add_paragraph()
    p.text = "Counselor logs contain qualitative text describing medical issues, financial panic, or student challenges."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)
    
    p = tf_doc.add_paragraph()
    p.text = "★ Unstructured and messy; cannot be fed directly to tabular models."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED

    # Arrow 1
    arrow1 = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.2), Inches(0.6), Inches(0.6))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = BORDER_COLOR
    arrow1.line.fill.background()

    # 2. Hexagon Process Block
    hex_shape = slide5.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(5.0), Inches(1.8), Inches(3.4), Inches(3.6))
    hex_shape.fill.solid()
    hex_shape.fill.fore_color.rgb = RGBColor(239, 246, 255) # Blue 50
    hex_shape.line.color.rgb = ACCENT_BLUE
    hex_shape.line.width = Pt(1.5)
    
    hex_box = slide5.shapes.add_textbox(Inches(5.3), Inches(2.0), Inches(2.8), Inches(3.2))
    tf_hex = hex_box.text_frame
    tf_hex.word_wrap = True
    p = tf_hex.paragraphs[0]
    p.text = "NLP PIPELINE"
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(14)
    
    p = tf_hex.add_paragraph()
    p.text = "1. TF-IDF Vectorization\nConverts counselor logs into word frequency metrics, focusing on key risk terms."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)
    
    p = tf_hex.add_paragraph()
    p.text = "2. Truncated SVD (LSA)\nReduces dimensions to dense semantic components, filtering noise."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER

    # Arrow 2
    arrow2 = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.6), Inches(3.2), Inches(0.6), Inches(0.6))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = BORDER_COLOR
    arrow2.line.fill.background()

    # 3. Output Circle
    circ_shape = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.4), Inches(1.8), Inches(3.2), Inches(3.6))
    circ_shape.fill.solid()
    circ_shape.fill.fore_color.rgb = RGBColor(240, 253, 250) # Emerald 50
    circ_shape.line.color.rgb = ACCENT_GREEN
    circ_shape.line.width = Pt(2)
    
    circ_box = slide5.shapes.add_textbox(Inches(9.6), Inches(2.1), Inches(2.8), Inches(3.0))
    tf_circ = circ_box.text_frame
    tf_circ.word_wrap = True
    p = tf_circ.paragraphs[0]
    p.text = "OUTPUT: DENSE SVD"
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(14)
    
    p = tf_circ.add_paragraph()
    p.text = "Dense numerical vectors integrating smoothly with tabular and time-series models, extracting semantic themes."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER

    # Summary box at bottom
    add_card(slide5, Inches(0.8), Inches(5.7), Inches(11.8), Inches(1.0), bg_color=CARD_BG, border_color=BORDER_COLOR)
    sum_box = slide5.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.4), Inches(0.8))
    tf_sum = sum_box.text_frame
    tf_sum.word_wrap = True
    p = tf_sum.paragraphs[0]
    p.text = "★ Performance Lift: Feature importance calculations prove that SVD embeddings from counselor logs rank as one of the strongest predictor variables in the final pipeline."
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    # ==========================================
    # SLIDE 6: Machine Learning Pipeline (Cards with Top Accent Strip)
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide6)
    add_slide_header(slide6, "Machine Learning Modeling Pipeline", "Modeling")

    # 3 Columns for 3 Base Models
    models_info = [
        {
            "name": "CatBoost",
            "desc": "Built-in support for categorical features. Uses ordered target statistics and symmetric trees.",
            "role": "Ensemble Weight: 70%",
            "color": ACCENT_BLUE
        },
        {
            "name": "LightGBM",
            "desc": "Fast, highly efficient histogram-based gradient booster. Uses leaf-wise tree growth.",
            "role": "Ensemble Weight: 15%",
            "color": ACCENT_GREEN
        },
        {
            "name": "XGBoost",
            "desc": "Extremely robust tree booster with built-in L1/L2 regularization to control complexity.",
            "role": "Ensemble Weight: 15%",
            "color": ACCENT_RED
        }
    ]

    for i, mod in enumerate(models_info):
        x = left_margin + i * (card_w + gap)
        
        # White card body
        add_card(slide6, x, Inches(2.0), card_w, Inches(4.0), bg_color=CARD_BG, border_color=BORDER_COLOR)
        
        # Color strip at top of card
        strip = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), card_w, Inches(0.2))
        strip.fill.solid()
        strip.fill.fore_color.rgb = mod["color"]
        strip.line.fill.background()
        
        box = slide6.shapes.add_textbox(x + Inches(0.2), Inches(2.4), card_w - Inches(0.4), Inches(3.4))
        tf_mod = box.text_frame
        tf_mod.word_wrap = True

        p = tf_mod.paragraphs[0]
        p.text = mod["name"].upper()
        p.font.name = FONT_TITLE
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = mod["color"]
        p.space_after = Pt(12)

        p = tf_mod.add_paragraph()
        p.text = mod["desc"]
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(20)

        p = tf_mod.add_paragraph()
        p.text = mod["role"]
        p.font.name = FONT_TITLE
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = mod["color"]

    # Add pipeline notes at the bottom
    pipe_box = slide6.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6))
    tf_pipe = pipe_box.text_frame
    tf_pipe.word_wrap = True
    p = tf_pipe.paragraphs[0]
    p.text = "Cross-Validation: Stratified 5-Fold to address class imbalance. | Primary optimization directed at Macro F1-Score."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 7: Ensemble & Threshold Optimization (Merging Flowchart)
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide7)
    add_slide_header(slide7, "Ensemble & Threshold Optimization", "Modeling")

    # Merging Diagram
    # Left Models
    models_list = [
        ("CatBoost (70%)", Inches(2.0)),
        ("LightGBM (15%)", Inches(3.1)),
        ("XGBoost (15%)", Inches(4.2))
    ]
    for name, y_pos in models_list:
        add_card(slide7, Inches(0.8), y_pos, Inches(2.6), Inches(0.9), bg_color=CARD_BG, border_color=BORDER_COLOR)
        box = slide7.shapes.add_textbox(Inches(0.9), y_pos + Inches(0.15), Inches(2.4), Inches(0.6))
        tf_m = box.text_frame
        tf_m.word_wrap = True
        p = tf_m.paragraphs[0]
        p.text = name
        p.font.name = FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        # Merging arrow from each
        arr = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.5), y_pos + Inches(0.25), Inches(0.4), Inches(0.4))
        arr.fill.solid()
        arr.fill.fore_color.rgb = BORDER_COLOR
        arr.line.fill.background()

    # Central Decision Diamond
    dec_diamond = slide7.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(4.1), Inches(2.4), Inches(3.2), Inches(2.3))
    dec_diamond.fill.solid()
    dec_diamond.fill.fore_color.rgb = RGBColor(239, 246, 255) # Blue 50
    dec_diamond.line.color.rgb = ACCENT_BLUE
    dec_diamond.line.width = Pt(1.5)
    
    dec_box = slide7.shapes.add_textbox(Inches(4.4), Inches(3.0), Inches(2.6), Inches(1.1))
    tf_dec = dec_box.text_frame
    tf_dec.word_wrap = True
    p = tf_dec.paragraphs[0]
    p.text = "Soft Voting Ensemble\n(Weighted Probs)"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Arrow to the right
    arr_right = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.5), Inches(3.35), Inches(0.5), Inches(0.4))
    arr_right.fill.solid()
    arr_right.fill.fore_color.rgb = BORDER_COLOR
    arr_right.line.fill.background()

    # Right Card: Threshold Optimization Detail
    add_card(slide7, Inches(8.2), Inches(1.8), Inches(4.3), Inches(4.4), bg_color=CARD_BG, border_color=ACCENT_GREEN, line_width=2)
    opt_box = slide7.shapes.add_textbox(Inches(8.4), Inches(2.0), Inches(3.9), Inches(4.0))
    tf_opt = opt_box.text_frame
    tf_opt.word_wrap = True

    p = tf_opt.paragraphs[0]
    p.text = "THRESHOLD OPTIMIZATION"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(12)

    p = tf_opt.add_paragraph()
    p.text = "Since the target metric is Macro F1, class imbalance severely penalizes poor predictions on minority classes."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

    p = tf_opt.add_paragraph()
    p.text = "By tuning threshold multipliers on validation folds, we scale minority probabilities:"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

    p = tf_opt.add_paragraph()
    p.text = "★ Class 1 Multiplier: 1.30\n★ Class 2 Multiplier: 1.30"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(10)

    p = tf_opt.add_paragraph()
    p.text = "Scaling shifts the decision boundaries, maximizing the recall of at-risk students and boosting the Macro F1."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 8: Model Performance (Bright styled Table)
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide8)
    add_slide_header(slide8, "Model Performance Comparison", "Results")

    # Left: Explanation card
    add_card(slide8, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.5), bg_color=CARD_BG, border_color=BORDER_COLOR)
    exp_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.1), Inches(4.1))
    tf_exp = exp_box.text_frame
    tf_exp.word_wrap = True

    p = tf_exp.paragraphs[0]
    p.text = "EVALUATION METRIC"
    p.font.name = FONT_TITLE
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(10)

    p = tf_exp.add_paragraph()
    p.text = "The primary evaluation metric is the Macro F1-Score, which computes F1 independently for each class and averages them."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

    p = tf_exp.add_paragraph()
    p.text = "Ensembling reduces error variances, while threshold optimization directly targets the metric imbalance penalty."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(24)

    p = tf_exp.add_paragraph()
    p.text = "Final Best Validation Score:"
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK

    p = tf_exp.add_paragraph()
    p.text = "0.7072 Macro F1"
    p.font.name = FONT_TITLE
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    # Add Bright Styled Table
    rows = 6
    cols = 2
    left = Inches(5.6)
    top = Inches(1.8)
    width = Inches(6.9)
    height = Inches(4.5)

    table_shape = slide8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(4.7)
    table.columns[1].width = Inches(2.2)

    headers = ["Model / Technique", "Validation Macro F1"]
    data = [
        ["CatBoost (Standalone)", "0.6991"],
        ["LightGBM (Standalone)", "0.6887"],
        ["XGBoost (Standalone)", "0.6927"],
        ["Three-Model Ensemble", "0.7017"],
        ["Ensemble + Threshold Optimization", "0.7072"]
    ]

    # Style Header Row
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
        p.font.name = FONT_TITLE
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Populate Data with clean light-theme styling
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            # Alternating background colors
            if row_idx == 4: # Highlight best ensemble row
                cell.fill.fore_color.rgb = RGBColor(209, 250, 229) # Emerald 100
            else:
                cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else RGBColor(241, 245, 249) # Slate 100
                
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
            p.font.name = FONT_BODY
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_DARK
            if row_idx == 4:
                p.font.bold = True
                p.font.color.rgb = ACCENT_GREEN

    # ==========================================
    # SLIDE 9: Key Insights & Challenges (Alternating Stripe Cards)
    # ==========================================
    slide9 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide9)
    add_slide_header(slide9, "Key Insights & Project Challenges", "Insights")

    # Left Column: Key Drivers (White card with Left Blue accent bar)
    add_card(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), bg_color=CARD_BG, border_color=BORDER_COLOR)
    # Stripe
    stripe_l = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(4.8))
    stripe_l.fill.solid()
    stripe_l.fill.fore_color.rgb = ACCENT_BLUE
    stripe_l.line.fill.background()

    ins_box = slide9.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_ins = ins_box.text_frame
    tf_ins.word_wrap = True

    p = tf_ins.paragraphs[0]
    p.text = "KEY DROPOUT DRIVERS"
    p.font.name = FONT_TITLE
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(14)

    drivers = [
        "Attendance Behavior: Strongest indicator of disengagement.",
        "CGPA Consistency & Trend: Sudden declines trigger warnings.",
        "Lifestyle Factors: Excessive screen time and long commute times.",
        "Socio-Economic Variables: Family income and parental support.",
        "Counsellor Insights: Highly discriminative when NLP-processed."
    ]
    for d in drivers:
        p = tf_ins.add_paragraph()
        parts = d.split(":")
        run1 = p.add_run()
        run1.text = "• " + parts[0] + ":"
        run1.font.bold = True
        run1.font.name = FONT_BODY
        run1.font.size = Pt(12)
        run1.font.color.rgb = TEXT_DARK
        
        run2 = p.add_run()
        run2.text = parts[1]
        run2.font.name = FONT_BODY
        run2.font.size = Pt(12)
        run2.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # Right Column: Challenges (White card with Left Red accent bar)
    add_card(slide9, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), bg_color=CARD_BG, border_color=BORDER_COLOR)
    # Stripe
    stripe_r = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(0.15), Inches(4.8))
    stripe_r.fill.solid()
    stripe_r.fill.fore_color.rgb = ACCENT_RED
    stripe_r.line.fill.background()

    chal_box = slide9.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_chal = chal_box.text_frame
    tf_chal.word_wrap = True

    p = tf_chal.paragraphs[0]
    p.text = "CHALLENGES OVERCOME"
    p.font.name = FONT_TITLE
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.space_after = Pt(14)

    chals = [
        "Messy Texts: Filtering repetitive boilerplate phrases from counselor reports.",
        "Feature Drift: Ensuring engineered academic/attendance metrics match across training and test splits.",
        "Class Disproportion: Overcoming poor validation scores on minority dropout classes via threshold optimization.",
        "Imbalanced F1 Optimization: Directing GBDT models toward Macro F1 rather than raw accuracy."
    ]
    for c in chals:
        p = tf_chal.add_paragraph()
        parts = c.split(":")
        run1 = p.add_run()
        run1.text = "• " + parts[0] + ":"
        run1.font.bold = True
        run1.font.name = FONT_BODY
        run1.font.size = Pt(12)
        run1.font.color.rgb = TEXT_DARK
        
        run2 = p.add_run()
        run2.text = parts[1]
        run2.font.name = FONT_BODY
        run2.font.size = Pt(12)
        run2.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 10: Future Scope & Takeaways (2x2 Grid)
    # ==========================================
    slide10 = prs.slides.add_slide(slide_layout)
    apply_bright_bg(slide10)
    add_slide_header(slide10, "Future Scope & Key Takeaways", "Conclusion")

    # 2x2 grid parameters
    grid_w = Inches(5.7)
    grid_h = Inches(2.1)
    grid_x = [Inches(0.8), Inches(6.8)]
    grid_y = [Inches(1.8), Inches(4.2)]

    grid_items = [
        {
            "col": 0, "row": 0,
            "title": "Feature Engineering Wins",
            "text": "Clean, domain-specific features (CGPA slopes, backlog trajectories) yield far higher performance than raw model parameter tuning.",
            "accent": ACCENT_BLUE
        },
        {
            "col": 1, "row": 0,
            "title": "Multimodal Synergy",
            "text": "Fusing socio-economics (tabular), temporal habits (attendance series), and emotional warnings (text) builds the most robust system.",
            "accent": ACCENT_BLUE
        },
        {
            "col": 0, "row": 1,
            "title": "Zero-Shot Text & SHAP (XAI)",
            "text": "Integrating Explainable AI (SHAP waterfall charts) and zero-shot NLP sentiment flags will make risks transparent and actionable.",
            "accent": ACCENT_GREEN
        },
        {
            "col": 1, "row": 1,
            "title": "Causal AI Interventions",
            "text": "Using Causal Inference models to suggest the exact treatment (tutoring, scholarship) that yields the maximum risk reduction.",
            "accent": ACCENT_GREEN
        }
    ]

    for item in grid_items:
        gx = grid_x[item["col"]]
        gy = grid_y[item["row"]]
        
        # White Grid Card
        add_card(slide10, gx, gy, grid_w, grid_h, bg_color=CARD_BG, border_color=BORDER_COLOR)
        
        # Tiny Accent Square at top left
        sq = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, gx + Inches(0.2), gy + Inches(0.2), Inches(0.2), Inches(0.2))
        sq.fill.solid()
        sq.fill.fore_color.rgb = item["accent"]
        sq.line.fill.background()
        
        box = slide10.shapes.add_textbox(gx + Inches(0.5), gy + Inches(0.1), grid_w - Inches(0.6), grid_h - Inches(0.2))
        tf_grid = box.text_frame
        tf_grid.word_wrap = True
        
        p = tf_grid.paragraphs[0]
        p.text = item["title"]
        p.font.name = FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = item["accent"]
        p.space_after = Pt(6)
        
        p = tf_grid.add_paragraph()
        p.text = item["text"]
        p.font.name = FONT_BODY
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED

    # Save to path
    out_path = r"C:\Users\hawka\Desktop\RetinaAI\Student_Dropout_Prediction_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == "__main__":
    create_deck()
